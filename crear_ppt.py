import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
from urllib.parse import urljoin

# --- CONFIGURACIÓN ---
URL_WEB = input("Introduce la URL completa: ")
NOMBRE_ARCHIVO = input("Nombre del archivo de salida: ")
if not NOMBRE_ARCHIVO.endswith(".pptx"): NOMBRE_ARCHIVO += ".pptx"

# Simular navegador
headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'}

print("Descargando web completa (esto puede tardar si hay muchas fotos)...")
try:
    response = requests.get(URL_WEB, headers=headers)
    response.raise_for_status()
except Exception as e:
    print(f"Error: {e}")
    exit()

soup = BeautifulSoup(response.text, 'html.parser')
prs = Presentation()

# --- VARIABLES DE CONTROL DE ESPACIO ---
# PowerPoint mide en "puntos" o "pulgadas". Vamos a usar un contador vertical "Y".
# Una diapositiva mide 7.5 pulgadas de alto. Usaremos hasta 6.5.
POSICION_Y_ACTUAL = Inches(1) # Empezamos un poco abajo
MARGEN_IZQUIERDO = Inches(0.5)
ANCHO_UTIL = Inches(9)
ALTO_MAXIMO = Inches(6.5)

# Función para crear nueva diapositiva cuando se llene la anterior
def nueva_diapositiva():
    global slide, POSICION_Y_ACTUAL
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # 6 = Blanco
    POSICION_Y_ACTUAL = Inches(0.5) # Reseteamos el cursor arriba
    return slide

# Creamos la primera
slide = nueva_diapositiva()

# --- FUNCIÓN: AÑADIR TEXTO ---
def agregar_texto(texto, tamano, negrita=False):
    global slide, POSICION_Y_ACTUAL
    
    # Estimación simple: Si nos pasamos de altura, cambio de slide
    if POSICION_Y_ACTUAL > ALTO_MAXIMO:
        slide = nueva_diapositiva()

    # Caja de texto
    alto_estimado = Inches(0.5) + (Inches(0.3) * (len(texto) / 100)) # Calculo "a ojo" de cuanto ocupará
    
    textbox = slide.shapes.add_textbox(MARGEN_IZQUIERDO, POSICION_Y_ACTUAL, ANCHO_UTIL, alto_estimado)
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = texto
    p.font.size = Pt(tamano)
    p.font.bold = negrita
    p.font.name = "Arial"
    
    # Actualizamos la posición del cursor para el siguiente elemento
    POSICION_Y_ACTUAL += alto_estimado + Inches(0.2) # Un poco de margen

# --- FUNCIÓN: AÑADIR IMAGEN ---
def agregar_imagen(src):
    global slide, POSICION_Y_ACTUAL
    
    if not src: return

    # Descargar
    full_url = urljoin(URL_WEB, src)
    try:
        r = requests.get(full_url, headers=headers, timeout=3)
        img_stream = BytesIO(r.content)
    except:
        return # Si falla, ignoramos

    # Si queda poco espacio, pasamos a la siguiente slide directamente
    if POSICION_Y_ACTUAL > (ALTO_MAXIMO - Inches(2)):
        slide = nueva_diapositiva()

    try:
        # Añadir imagen
        pic = slide.shapes.add_picture(img_stream, MARGEN_IZQUIERDO, POSICION_Y_ACTUAL, width=Inches(4))
        
        # Mover cursor (bajamos lo que mida la imagen)
        POSICION_Y_ACTUAL += pic.height + Inches(0.2)
    except:
        pass

# --- BUCLE PRINCIPAL (EL CEREBRO DEL SCRIPT) ---
# Buscamos TODOS los elementos relevantes en orden
contenido = soup.find_all(['h1', 'h2', 'h3', 'p', 'img', 'li'])

print(f"Procesando {len(contenido)} elementos encontrados...")

for elemento in contenido:
    
    # 1. ES UN TÍTULO
    if elemento.name in ['h1', 'h2', 'h3']:
        texto = elemento.get_text().strip()
        if texto:
            # Los títulos siempre fuerzan un poco de aire o nueva slide si estamos muy abajo
            if POSICION_Y_ACTUAL > Inches(5): 
                slide = nueva_diapositiva()
            agregar_texto(texto, 24, True) # Grande y negrita

    # 2. ES TEXTO NORMAL (Párrafos o Listas)
    elif elemento.name in ['p', 'li']:
        texto = elemento.get_text().strip()
        if texto:
            agregar_texto(texto, 12, False) # Normal

    # 3. ES UNA IMAGEN
    elif elemento.name == 'img':
        src = elemento.get('src')
        if src:
            agregar_imagen(src)

# --- GUARDAR ---
try:
    prs.save(NOMBRE_ARCHIVO)
    print(f"✅ ¡TERMINADO! Se han generado las diapositivas con todo el contenido.")
    print(f"Archivo: {NOMBRE_ARCHIVO}")
except PermissionError:
    print("❌ Error: Cierra el PowerPoint antes de generarlo.")